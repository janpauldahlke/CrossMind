"""Generate CrossMind pitch deck (7 slides) for AI BEAVERS founder hackathon."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(
    os.path.expanduser("~"),
    ".cursor/projects/Users-jandahlke-dev-hagbards-stuff-crossmind/assets",
)

IMG_SEALED = os.path.join(
    ASSETS_DIR, "2_sealed_with_pass-1d582270-e737-4d46-9f4e-3c3e75c8366f.png"
)
IMG_SEALED_WRONG = os.path.join(
    ASSETS_DIR, "3_sealed_incorrect_pass-65f9d005-2d82-4181-b521-64e19af025ea.png"
)
IMG_HELIX = os.path.join(
    ASSETS_DIR, "4_helix_routing_encrypt-c0c5322f-4c4c-451d-8ef5-06407fdf1fa8.png"
)
IMG_HELIX_WRONG = os.path.join(
    ASSETS_DIR, "5_helix_routing_incorrect_key-9dcf372d-cdb3-4c5b-8c94-c5e5f35a279f.png"
)
IMG_INITIAL = os.path.join(
    ASSETS_DIR, "1_initial_state-7c9e8b79-2041-4e92-b799-2a1511f5b9b4.png"
)

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x99, 0x99, 0xAA)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide):
    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(12)
    height = Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DEMO \u2014 Research prototype only. Not medical advice, diagnosis, or treatment."
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT


def add_title(slide, text, top=Inches(0.6), left=Inches(0.8), size=Pt(36)):
    width = Inches(11)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = WHITE
    return txBox


def add_subtitle(slide, text, top=Inches(1.6), left=Inches(0.8)):
    width = Inches(11)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    return txBox


def add_bullet_block(slide, bullets, top=Inches(2.6), left=Inches(0.8), width=Inches(11)):
    height = Inches(4.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (main, sub) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(14)
        run = p.add_run()
        run.text = main
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if sub:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            run2 = p2.add_run()
            run2.text = "    " + sub
            run2.font.size = Pt(16)
            run2.font.color.rgb = LIGHT_GRAY
    return txBox


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ─── SLIDE 1: Problem + broad intro ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "Every large organization wants AI.\nNone of them want to export their data.", size=Pt(34))
    add_subtitle(slide, "CrossMind \u2014 AI BEAVERS Founder Hackathon \u00b7 June 6, 2026", top=Inches(2.0))
    add_bullet_block(slide, [
        ("\u2022  The problem is universal", "Healthcare, finance, logistics, legal, industrial \u2014 sensitive text cannot leave the local network."),
        ("\u2022  The buyer", "Every organization with a large specialist model at HQ and smaller models at the edge."),
        ("\u2022  Today\u2019s choice", "Ban AI entirely \u2014 or export full text to a cloud LLM and accept the audit nightmare."),
        ("\u2022  For this hackathon", "We demonstrate the approach in healthcare: clinic \u2194 hospital."),
    ], top=Inches(3.0))
    add_footer(slide)

    # ─── SLIDE 2: Solution ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "CrossMind \u2014 split inference across two models")
    add_bullet_block(slide, [
        ("\u2022  Local model: Qwen 2.5 (clinic)", "Encodes patient text locally. The prompt string never leaves the clinic."),
        ("\u2022  Remote model: Llama 3.1 (hospital)", "Runs the specialist output layer on aligned vectors \u2014 not on your sentence."),
        ("\u2022  Linear alignment", "A learned matrix W* maps Qwen\u2019s hidden space into Llama\u2019s basis (Gorbett & Jana, 2025)."),
        ("\u2022  Trained on medical Q&A", "Alignment trained on clinical question-answer pairs for domain accuracy."),
    ], top=Inches(2.2))

    # Architecture diagram as text
    left = Inches(0.8)
    top = Inches(5.8)
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinic (Qwen)  \u2192  h_B  \u2192  h_aligned = h_B \u00d7 W* + b*  \u2500\u2500 vectors \u2500\u2500\u25b6  Hospital (Llama) \u2192 LM head \u2192 token"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = "Courier New"
    add_footer(slide)

    # ─── SLIDE 3: Privacy Layer 1 \u2014 Sealed ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "Privacy Layer 1 \u2014 Sealed (wire obfuscation)")
    add_bullet_block(slide, [
        ("\u2022  Vector-by-vector transmission", "Each hidden state is sent individually \u2014 not the full document."),
        ("\u2022  Obfuscation with shared passphrase (rotation R)", "h_enc = h_aligned \u00d7 R  \u2014  Specialist only sees rotated vectors, never raw text."),
        ("\u2022  Wrong passphrase = garbled output", "Without the matching key, intercepted vectors are meaningless noise."),
        ("\u2022  Protects against", "Passive eavesdropper on the network link."),
    ], top=Inches(2.2), width=Inches(5.8))

    add_image_safe(slide, IMG_SEALED, Inches(6.8), Inches(2.0), width=Inches(6.0))
    add_footer(slide)

    # ─── SLIDE 4: Privacy Layer 2 \u2014 HELIX ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "Privacy Layer 2 \u2014 HELIX (homomorphic encryption)")
    add_bullet_block(slide, [
        ("\u2022  CKKS homomorphic encryption", "The hospital computes on ciphertext \u2014 never sees the plaintext vector."),
        ("\u2022  Department routing (5 classes)", "Cardiology, Neurology, Oncology, Orthopedics, General Medicine."),
        ("\u2022  Hospital never sees the department label", "Only the clinic can decrypt the routing result. Truly encrypted compute."),
        ("\u2022  ~3.4s per routing request", "Practical for classification heads. Not yet viable for full-vocab generation."),
    ], top=Inches(2.2), width=Inches(5.8))

    add_image_safe(slide, IMG_HELIX, Inches(6.8), Inches(2.0), width=Inches(6.0))
    add_footer(slide)

    # ─── SLIDE 5: Privacy overview \u2014 what it does / doesn't ─────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "What our privacy layers do \u2014 and what they don\u2019t")

    # "Does" column
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "\u2705  What it DOES"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = ACCENT_GREEN

    does_items = [
        "Prompt text stays on the local machine",
        "Wire tap sees only rotated/encrypted vectors",
        "HELIX: hospital computes without seeing plaintext",
        "Department label decrypted only by clinic",
        "Data minimization \u2014 no full-text export",
    ]
    for item in does_items:
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE

    # "Doesn't" column
    txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = "\u26a0\ufe0f  What it does NOT (yet)"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = ORANGE

    doesnt_items = [
        "Full encrypted generation (server decrypts for Sealed gen)",
        "HIPAA/GDPR compliance certification",
        "Guarantee hospital learns nothing from vectors",
        "Clinical-grade diagnosis or medical advice",
        "Protection from a curious hub (for generation)",
    ]
    for item in doesnt_items:
        p = tf2.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(16)
        run.font.color.rgb = LIGHT_GRAY

    add_footer(slide)

    # ─── SLIDE 6: Results + screenshots ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "Live demo \u2014 Clinic \u2194 Hospital in action", size=Pt(30))

    # Metrics callout
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Alignment cosine: ~0.83   |   Routing accuracy: ~83% val   |   HELIX compute: ~3.4s   |   Sealed: real-time streaming"
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_GREEN

    add_image_safe(slide, IMG_INITIAL, Inches(0.3), Inches(2.5), width=Inches(6.3))
    add_image_safe(slide, IMG_HELIX, Inches(6.8), Inches(2.5), width=Inches(6.3))

    # Labels
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(6.2), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u2191 Sealed mode: clinic encodes, hospital decodes from vectors"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(6.2), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u2191 HELIX mode: encrypted routing, clinic decrypts department"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide)

    # ─── SLIDE 7: Phase 2 + Thank you ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_title(slide, "Phase 2 \u2014 Beyond text queries")
    add_bullet_block(slide, [
        ("\u2022  Upload documents (PDF, TXT)", "Send encrypted document vectors for specialist analysis without exposing content."),
        ("\u2022  Upload images & reports", "Medical imaging, lab reports \u2014 encoded locally, routed or analyzed remotely."),
        ("\u2022  Multi-modal encrypted routing", "Same HELIX pattern extended to richer inputs for better clinical results."),
        ("\u2022  The approach is domain-agnostic", "Finance, logistics, legal, industrial \u2014 wherever sensitive data meets central AI."),
    ], top=Inches(2.2))

    # Thank you + reference
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank you"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT_PURPLE
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "Reference: Gorbett, T., & Jana, S. (2025). Characterizing Linear Alignment Across Language Models. arXiv:2603.18908"
    run2.font.size = Pt(12)
    run2.font.color.rgb = MUTED
    p3 = tf.add_paragraph()
    p3.space_before = Pt(6)
    run3 = p3.add_run()
    run3.text = "AI BEAVERS Founder Hackathon \u00b7 House of AI Hamburg \u00b7 June 6, 2026"
    run3.font.size = Pt(12)
    run3.font.color.rgb = MUTED

    add_footer(slide)

    # Save
    out_path = os.path.join(SCRIPT_DIR, "CrossMind_Pitch_Deck.pptx")
    prs.save(out_path)
    print(f"Deck saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_deck()
